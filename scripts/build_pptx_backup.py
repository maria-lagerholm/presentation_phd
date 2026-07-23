from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from pptx_helpers import (
    NAVY, GREEN, GOLD, RUST, GRAY, LIGHT, WHITE,
    ROOT,
    add_textbox, write_text, add_para, add_title, add_kicker,
    add_notes, add_picture, add_circular_picture, add_zoomable_picture,
    add_footer, set_fade_transition, blank_slide, set_slide_size,
)

TEAL = RGBColor(0x20, 0xB2, 0xAA)
MID_BLUE = RGBColor(0x3D, 0x6A, 0x99)
CARD = RGBColor(0xF7, 0xF7, 0xF5)
GREEN_TINT = RGBColor(0xEC, 0xF1, 0xEA)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_rounded(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    return shape


def fill_box(box, lines, align=PP_ALIGN.LEFT):
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, size, bold, color in lines:
        if first:
            write_text(tf, text, size=size, bold=bold, color=color, align=align)
            first = False
        else:
            p = add_para(tf, text, size=size, bold=bold, color=color)
            p.alignment = align


def caption(slide, text, left, top, width):
    box = add_textbox(slide, left, top, width, Inches(0.3))
    write_text(box.text_frame, text, size=10, color=LIGHT, align=PP_ALIGN.CENTER)


def data_table(slide, rows, left, top, width, height, font_size=12):
    n_rows = len(rows)
    n_cols = len(rows[0])
    gf = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for r in range(n_rows):
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            cell.text = str(rows[r][c])
            cell.vertical_anchor = 3
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                run.font.color.rgb = GRAY
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return gf


def slide_title(prs):
    slide = blank_slide(prs)
    add_picture(slide, "ub.png", Inches(0.7), Inches(0.7), height=Inches(2.0))
    add_picture(slide, "agrif.png", Inches(10.9), Inches(0.7), height=Inches(2.0))
    box = add_textbox(slide, Inches(2.6), Inches(1.35), Inches(8.1), Inches(0.6))
    write_text(box.text_frame, "University of Belgrade  ·  Faculty of Agriculture", size=17, color=GRAY, align=PP_ALIGN.CENTER)

    title = add_textbox(slide, Inches(1.0), Inches(2.95), Inches(11.3), Inches(2.0))
    write_text(
        title.text_frame,
        "Chemical characterisation of essential oils and machine learning analysis of molecular substructures associated with antimicrobial activity against Salmonella Typhimurium for food industry applications",
        size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )
    label = add_textbox(slide, Inches(1.0), Inches(5.05), Inches(11.3), Inches(0.4))
    write_text(label.text_frame, "Doctoral Defense", size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    author = add_textbox(slide, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2))
    fill_box(author, [
        ("Student: Maria Lagerholm", 16, False, NAVY),
        ("Mentors: Prof Dr Viktor Nedović, Prof Dr Mirjana Pešić", 15, False, GRAY),
        ("2026", 16, False, GREEN),
    ], align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "Good morning everyone. Today I'll be presenting my doctoral research on essential oils and how we can use machine learning to get a better understanding of what drives their antimicrobial activity, specifically against Salmonella.\n\nThis work combines lab experiments with computational analysis. It was done at the Faculty of Agriculture, University of Belgrade, in collaboration with other faculties and universities.")


def slide_acknowledgments(prs):
    slide = blank_slide(prs)
    add_title(slide, "Acknowledgments", size=30)
    sub = add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.35))
    write_text(sub.text_frame, "With gratitude to all collaborators who contributed to this research", size=14, italic=True, color=GRAY)

    people = [
        [("people/viktor.jpg", "Prof Dr Viktor Nedović", "Faculty of Agriculture, UB"),
         ("people/mirjana.jpg", "Prof Dr Mirjana Pešić", "Faculty of Agriculture, UB")],
        [("people/anaS.jpg", "Dr Ana Salević", "Faculty of Agriculture, UB"),
         ("people/anaK.jpg", "Dr Ana Kalušević", "Institute of Meat Hygiene"),
         ("people/steva.jpg", "Dr Steva Lević", "Faculty of Agriculture, UB"),
         ("people/milena.jpg", "Dr Milena Pantić", "Faculty of Agriculture, UB"),
         ("people/miomir.jpeg", "Prof Dr Miomir Nikšić", "Faculty of Agriculture, UB")],
        [("people/milos.jpg", "Dr Miloš Jovanović", "Faculty of Organizational Sciences"),
         ("people/marija-kuzmanovic.jpeg", "Dr Marija Kuzmanović", "Faculty of Organizational Sciences"),
         ("people/dejan.jpg", "Dr Dejan Pljevljakušić", "Institute Dr Josif Pančić"),
         ("people/Katarina-Savikin.jpg", "Dr Katarina Šavikin", "Institute Dr Josif Pančić"),
         ("people/Natasa-Milosavljevic.jpg", "Prof Dr Nataša Milosavljević", "Faculty of Agriculture, UB")],
        [("people/ivan.jpg", "Prof Dr Ivan Mijaković", "Chalmers & DTU"),
         ("people/nielsen.webp", "Prof Dr Jens Nielsen", "BioInnovation Institute & Chalmers"),
         ("people/aleksej.jpg", "Prof Dr Aleksej Zelezniak", "Chalmers & Vilnius"),
         ("people/jan_zrimec.jpeg", "Dr Jan Zrimec", "NIB Slovenia & Chalmers"),
         ("people/filip.jpeg", "Filip Buric", "Chalmers University")],
    ]

    name_h = 0.5
    gap = 0.16
    y = 1.3
    for row in people:
        n = len(row)
        card_w = 2.2 if n > 2 else 3.2
        diameter = 0.95 if n <= 2 else 0.8
        col_gap = 0.15
        total = n * card_w + (n - 1) * col_gap
        x0 = (13.333 - total) / 2
        for i, (img, name, aff) in enumerate(row):
            x = x0 + i * (card_w + col_gap)
            add_circular_picture(slide, img, Inches(x + (card_w - diameter) / 2), Inches(y),
                                  Inches(diameter), ring=GOLD, ring_pt=1.5)
            nb = add_textbox(slide, Inches(x), Inches(y + diameter + 0.03), Inches(card_w), Inches(name_h))
            fill_box(nb, [(name, 9, True, NAVY), (aff, 8, False, LIGHT)], align=PP_ALIGN.CENTER)
        y += diameter + name_h + gap

    add_footer(slide)
    add_notes(slide, "Before we start, I want to thank everyone who helped make this work possible.\n\nMy mentors at the Faculty of Agriculture for guiding me through this journey. The experimental team here in Belgrade who helped with the laboratory work. Our collaborators at the Faculty of Organizational Sciences for the computational side. The Institute for Medicinal Plants team for their botanical expertise. And my supervisors during international mobility at Chalmers University and other institutions, who brought computational biology knowledge to this project.")


def slide_roadmap(prs):
    slide = blank_slide(prs)
    add_title(slide, "Research Roadmap")
    add_picture(slide, "plants/botanical-lab-work-cut.png", Inches(1.2), Inches(1.6), height=Inches(2.1))
    add_picture(slide, "plants/botanical-ml-analysis-cut.png", Inches(8.6), Inches(1.6), height=Inches(2.1))

    left = add_textbox(slide, Inches(0.7), Inches(3.9), Inches(5.0), Inches(2.6))
    fill_box(left, [
        ("PHASE I: FOUNDATION", 12, True, GREEN),
        ("Experimental Lab Work", 20, True, NAVY),
        ("GC MS profiling and bioactivity validation of 4 commercial EOs; development of calcium alginate delivery systems.", 14, False, GRAY),
        ("Outcome: Encapsulation Strategy", 13, True, GREEN),
    ])

    right = add_textbox(slide, Inches(7.7), Inches(3.9), Inches(5.0), Inches(2.6))
    fill_box(right, [
        ("PHASE II: DATA DRIVEN ANALYSIS", 12, True, NAVY),
        ("Machine Learning & Meta Analysis", 20, True, NAVY),
        ("Decoding molecular drivers of antimicrobial activity through large scale data mining.", 14, False, GRAY),
        ("Outcome: Molecular Insight", 13, True, NAVY),
    ])
    add_footer(slide)
    add_notes(slide, "Here's how today's presentation is organized.\n\nFirst, Phase I: the experimental work. We characterized commercial essential oils and suggested a way to preserve and deliver them using encapsulation in biodegradable material.\n\nThen Phase II: the machine learning part. We collected and analysed data to figure out what molecular structures are potentially involved in the antimicrobial activity of plant essential oils.\n\nI'll walk you through the experiments, explain the gap we found that led us to the computational work, and then show you what the machine learning revealed.")


def slide_why_hub(prs):
    slide = blank_slide(prs)
    title = add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12.1), Inches(1.0))
    write_text(title.text_frame, "Why?", size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    items = [
        ("Why did we focus on Salmonella?", "Major foodborne pathogen · food safety model · survives in food · inflammation · test bed for new antimicrobials", RUST),
        ("Why study essential oils in 2026?", "Plant chemical innovation · complex bioactive mixtures · ML & CRISPR · engineered microbes · beyond extraction", GREEN),
        ("Why machine learning?", "Hundreds of compounds · impractical individual testing · patterns across many oils", NAVY),
    ]
    y = 2.1
    for title_t, body, color in items:
        add_rounded(slide, Inches(1.4), Inches(y), Inches(10.5), Inches(1.25), CARD)
        add_rect(slide, Inches(1.4), Inches(y), Inches(0.14), Inches(1.25), color)
        box = add_textbox(slide, Inches(1.8), Inches(y + 0.12), Inches(9.9), Inches(1.05))
        fill_box(box, [(title_t, 18, True, NAVY), (body, 13, False, GRAY)])
        y += 1.5
    add_footer(slide)
    add_notes(slide, "Why study essential oils in 2026? Because plants have already done a lot of chemical innovation for us. Essential oils are complex mixtures of bioactive molecules derived from plants, and today we have tools such as machine learning, synthetic biology, and CRISPR based genome editing. If we can understand which plant molecules are useful, we may later be able to produce them in engineered microbes instead of relying only on extraction from plants.\n\nWhy did we focus on Salmonella? Because it is a major foodborne pathogen and an important model for food safety research. It can survive in foods, cause intestinal infection, trigger strong inflammation, and develop resistance to antibiotics. Historically, Salmonella infections caused serious outbreaks. Today, sanitation and hygiene allow us to control it much better, but it remains a useful target for testing new antimicrobial strategies.\n\nWhy machine learning? There are hundreds of compounds in essential oils. Testing them one by one is impractical. Machine learning allows us to find patterns across many oils at once.")


def slide_lab_hub(prs):
    slide = blank_slide(prs)
    add_kicker(slide, "Phase I", GREEN)
    add_title(slide, "Lab Experiments", top=0.45)

    plants = [
        ("plants/botanical-teatree-cut.png", "Tea tree", "terpinen-4-ol"),
        ("plants/botanical-lavender-cut.png", "Lavender", "linalyl acetate"),
        ("plants/botanical-bergamot-cut.png", "Bergamot", "linalyl acetate"),
        ("plants/botanical-peppermint-cut.png", "Peppermint", "isomenthol"),
    ]
    for i, (img, name, mol) in enumerate(plants):
        x = 1.2 + i * 2.9
        add_picture(slide, img, Inches(x), Inches(1.25), height=Inches(1.35))
        nb = add_textbox(slide, Inches(x - 0.2), Inches(2.65), Inches(2.4), Inches(0.6))
        fill_box(nb, [(name, 12, True, NAVY), (mol, 11, False, GREEN)], align=PP_ALIGN.CENTER)

    cards = [
        ("01 · Chemotyping", "GC MS", "Four oils · >98% volatiles", "gc-ms.png", GREEN),
        ("02 · Bioactivity", "MIC & antioxidants", "S. Typhimurium · DPPH / ABTS", "antimicrob.png", NAVY),
        ("03 · Delivery", "Alginate beads", "2.0 to 1.4 mm", "beads-shrinkage.png", GOLD),
    ]
    for i, (kicker, title, body, img, color) in enumerate(cards):
        x = 0.8 + i * 4.1
        add_rounded(slide, Inches(x), Inches(3.4), Inches(3.8), Inches(2.25), CARD)
        add_rect(slide, Inches(x), Inches(3.4), Inches(3.8), Inches(0.08), color)
        add_picture(slide, img, Inches(x + 1.15), Inches(3.6), height=Inches(0.7))
        box = add_textbox(slide, Inches(x + 0.2), Inches(4.4), Inches(3.4), Inches(1.15))
        fill_box(box, [(kicker, 11, True, color), (title, 16, True, NAVY), (body, 12, False, GRAY)])

    add_rounded(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1.05), RGBColor(0xF7, 0xE9, 0xE3))
    gap = add_textbox(slide, Inches(1.1), Inches(5.95), Inches(11.1), Inches(0.9))
    fill_box(gap, [
        ("RESEARCH GAP", 11, True, RUST),
        ("Absence of a computational way to map active chemical motifs to plants, so we no longer need to test compounds one by one in the lab.", 13, False, NAVY),
    ])
    add_footer(slide)
    add_notes(slide, "We started in the lab, using GC-MS to profile four essential oils: tea tree, lavender, bergamot, and peppermint. In each one we identified over 98% of the volatiles. Tea tree was mainly terpinen-4-ol (45%), lavender and bergamot were rich in linalool and linalyl acetate, and peppermint was mostly isomenthol (49%).\n\nNext we tested antimicrobial activity against Salmonella Typhimurium ATCC 14028. Lavender and bergamot inhibited growth at 5 µg/mL, while tea tree and peppermint required 10 µg/mL. For antioxidant activity we used DPPH and ABTS: peppermint was the strongest antioxidant (23 mmol TE/L by DPPH), tea tree the weakest (8 mmol TE/L). Activity remained measurable after 12 months of storage.\n\nThe problem is that raw oils evaporate and are hard to dose. So we encapsulated them into calcium alginate beads by electrostatic extrusion, a solvent free method. The wet beads are unstable, so we freeze dried them. Empty beads shrank from about 2.0 mm to 1.4 mm, roughly 30%, and stayed approximately spherical.\n\nThe research gap is not that bioactive compounds were unknown. For many years we learned which molecules work mainly by empirical laboratory testing, compound by compound and oil by oil. What was missing is a computational way to map chemical motifs linked to activity back to plants, so that this knowledge can be reached more systematically and with less labour. Now that the technology allows it, machine learning can help us leave that purely empirical path behind and prioritise what to test next. That is why we move to the computational part of the work.")


def slide_gcms(prs):
    slide = blank_slide(prs)
    add_kicker(slide, "01 · Chemotyping", GREEN)
    add_title(slide, "GC MS profiling", top=0.45)
    left = add_textbox(slide, Inches(0.7), Inches(1.6), Inches(6.5), Inches(4.5))
    fill_box(left, [
        ("Four commercial oils · >98% of volatiles identified", 15, False, GRAY),
        ("", 8, False, GRAY),
        ("Tea tree : Terpinen-4-ol · 45%", 17, True, GREEN),
        ("Lavender : Linalyl acetate · 42%", 17, True, NAVY),
        ("Bergamot : Linalyl acetate · 58%", 17, True, GOLD),
        ("Peppermint : Isomenthol · 49%", 17, True, TEAL),
    ])
    add_zoomable_picture(prs, slide, "gc-ms.png", Inches(8.0), Inches(1.35), height=Inches(4.9))
    add_footer(slide)
    add_notes(slide, "We started in the lab, using GC-MS to profile four essential oils: tea tree, lavender, bergamot, and peppermint. In each one we identified over 98% of the volatiles. Tea tree was mainly terpinen-4-ol (45%), lavender and bergamot were rich in linalool and linalyl acetate, and peppermint was mostly isomenthol (49%).")


def slide_bioactivity(prs):
    slide = blank_slide(prs)
    add_kicker(slide, "02 · Bioactivity", NAVY)
    add_title(slide, "S. Typhimurium ATCC 14028", top=0.45)

    add_rounded(slide, Inches(0.6), Inches(1.5), Inches(5.9), Inches(4.9), CARD)
    h1 = add_textbox(slide, Inches(0.95), Inches(1.65), Inches(5.2), Inches(0.5))
    write_text(h1.text_frame, "Antimicrobial · MIC", size=17, bold=True, color=NAVY)
    ls = add_textbox(slide, Inches(1.0), Inches(2.5), Inches(2.6), Inches(1.6))
    fill_box(ls, [("5", 46, True, NAVY), ("µg/mL", 15, False, GRAY), ("Lavender · Bergamot", 13, True, GRAY)], align=PP_ALIGN.CENTER)
    rs = add_textbox(slide, Inches(3.6), Inches(2.5), Inches(2.6), Inches(1.6))
    fill_box(rs, [("10", 46, True, GREEN), ("µg/mL", 15, False, GRAY), ("Tea tree · Peppermint", 13, True, GRAY)], align=PP_ALIGN.CENTER)
    add_zoomable_picture(prs, slide, "antimicrob.png", Inches(1.7), Inches(4.5), width=Inches(3.7))
    caption(slide, "MIC table 4.5", Inches(1.0), Inches(6.0), Inches(5.1))

    add_rounded(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(4.9), CARD)
    h2 = add_textbox(slide, Inches(7.15), Inches(1.65), Inches(5.2), Inches(0.5))
    write_text(h2.text_frame, "Antioxidant · mmol TE/L", size=17, bold=True, color=NAVY)
    data_table(slide, [
        ["Oil", "DPPH", "ABTS"],
        ["Tea tree", "8", "169"],
        ["Lavender", "19", "49"],
        ["Bergamot", "17", "99"],
        ["Peppermint", "23", "191"],
    ], Inches(7.3), Inches(2.3), Inches(5.0), Inches(2.0), font_size=12)
    add_zoomable_picture(prs, slide, "antiox.png", Inches(8.35), Inches(4.45), height=Inches(1.25))
    caption(slide, "Table 4.4", Inches(6.9), Inches(5.75), Inches(5.7))
    note = add_textbox(slide, Inches(6.9), Inches(6.05), Inches(5.7), Inches(0.4))
    write_text(note.text_frame, "Fresh oils · still measurable after 12 months", size=12, color=LIGHT, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_notes(slide, "Next we tested antimicrobial activity against Salmonella Typhimurium ATCC 14028. Lavender and bergamot inhibited growth at 5 µg/mL, while tea tree and peppermint required 10 µg/mL. For antioxidant activity we used DPPH and ABTS: peppermint was the strongest antioxidant (23 mmol TE/L by DPPH), tea tree the weakest (8 mmol TE/L). Activity remained measurable after 12 months of storage.")


def slide_encap(prs):
    slide = blank_slide(prs)
    add_kicker(slide, "03 · Delivery", GOLD)
    add_title(slide, "Calcium alginate encapsulation", top=0.45)

    add_zoomable_picture(prs, slide, "beads-shrinkage.png", Inches(1.6), Inches(1.4), width=Inches(10.1))
    caption(slide, "Before (top) and after (bottom) freeze drying", Inches(1.6), Inches(4.5), Inches(10.1))

    left = add_textbox(slide, Inches(0.8), Inches(5.0), Inches(7.2), Inches(1.7))
    fill_box(left, [
        ("Electrostatic extrusion · solvent free · freeze dried for stability", 15, False, GRAY),
        ("2.0 mm wet  →  1.4 mm dry  (~30% shrinkage)", 18, True, GOLD),
        ("Empty beads stayed approximately spherical after freeze drying.", 13, False, GRAY),
    ])
    add_zoomable_picture(prs, slide, "extrusion.png", Inches(9.8), Inches(4.9), height=Inches(1.75))
    caption(slide, "Electrostatic extrusion", Inches(8.6), Inches(6.7), Inches(4.2))
    add_footer(slide)
    add_notes(slide, "The problem is that raw oils evaporate and are hard to dose. So we encapsulated them into calcium alginate beads by electrostatic extrusion, a solvent free method. The wet beads are unstable, so we freeze dried them. Empty beads shrank from about 2.0 mm to 1.4 mm, roughly 30%, and stayed approximately spherical.\n\nThe research gap is not that bioactive compounds were unknown. For many years we learned which molecules work mainly by empirical laboratory testing, compound by compound and oil by oil. What was missing is a computational way to map chemical motifs linked to activity back to plants, so that this knowledge can be reached more systematically and with less labour. Now that the technology allows it, machine learning can help us leave that purely empirical path behind and prioritise what to test next. That is why we move to the computational part of the work.")


def slide_ml_pipeline(prs):
    slide = blank_slide(prs)
    add_kicker(slide, "Phase II", NAVY)
    add_title(slide, "Machine Learning", top=0.45)
    sub = add_textbox(slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.35))
    write_text(sub.text_frame, "A four step pipeline from literature oils to interpretable motifs", size=14, color=LIGHT)

    steps = [
        ("1", "Collect data", "Published EO studies · composition + anti Salmonella labels · majors >10%", "171 oils · 87 / 84", NAVY),
        ("2", "Encode molecules", "Names → SMILES → binary Morgan fingerprints", "bits · molecular features", MID_BLUE),
        ("3", "Model & rank", "L1 logistic regression + permutation feature importance", "PFI · key bits", GREEN),
        ("4", "Check literature", "Compare selected motifs with known EO bioactivity", "validate", GOLD),
    ]
    for i, (num, title, body, stat, color) in enumerate(steps):
        x = 0.5 + i * 3.2
        add_rounded(slide, Inches(x), Inches(1.55), Inches(2.95), Inches(3.55), CARD)
        add_rect(slide, Inches(x), Inches(1.55), Inches(2.95), Inches(0.1), color)
        nbox = add_textbox(slide, Inches(x + 1.05), Inches(1.8), Inches(0.85), Inches(0.5))
        write_text(nbox.text_frame, num, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        tbox = add_textbox(slide, Inches(x + 0.2), Inches(2.45), Inches(2.6), Inches(2.5))
        fill_box(tbox, [(title, 15, True, NAVY), (body, 12, False, GRAY), ("", 8, False, GRAY), (stat, 14, True, color)])

    add_rounded(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.35), GREEN_TINT, line_color=GREEN)
    outcome = add_textbox(slide, Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.1))
    fill_box(outcome, [
        ("OUTCOME", 12, True, GREEN),
        ("A reusable pipeline: the same approach can guide which compounds to test next in the lab.", 17, True, NAVY),
    ])
    add_footer(slide)
    add_notes(slide, "Here we built a four step pipeline.\n\nStep 1: We collected data on 171 essential oils from published studies, including their composition and whether they were classified as active against Salmonella. The set was nearly balanced: 87 active and 84 inactive oils. Only major constituents above 10% were included.\n\nStep 2: We converted the chemical names into Morgan fingerprints. These are binary codes that represent molecular features. Each chemical becomes a specific pattern of ones and zeros, where a one indicates the presence of a particular feature.\n\nStep 3: We trained an L1-regularised logistic regression model and then used permutation feature importance to determine which molecular features mattered most.\n\nStep 4: We compared our findings with what is already known from the essential oil literature.\n\nThe main outcome is the machine learning pipeline we developed. The same approach can be used to identify molecular features linked to other biological activities and to decide more efficiently what should be tested next in the laboratory.")


def slide_findings(prs):
    slide = blank_slide(prs)
    add_title(slide, "Machine Learning Insights")

    ftitle = add_textbox(slide, Inches(0.6), Inches(1.15), Inches(3.6), Inches(0.35))
    write_text(ftitle.text_frame, "FEATURE SELECTION", size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    stages = [
        ("2048 Morgan bits", 3.4, NAVY),
        ("682 present in data", 2.75, MID_BLUE),
        ("337 uncorrelated", 2.1, GREEN),
        ("10 key bits", 1.45, GOLD),
    ]
    center_x = 2.4
    for i, (label, w, color) in enumerate(stages):
        top = 1.6 + i * 0.6
        bar = add_rounded(slide, Inches(center_x - w / 2), Inches(top), Inches(w), Inches(0.5), color)
        tf = bar.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_text(tf, label, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    note = add_textbox(slide, Inches(0.6), Inches(4.05), Inches(3.6), Inches(0.4))
    write_text(note.text_frame, "PFI threshold: ROC AUC drop of at least 0.01", size=11, color=LIGHT, align=PP_ALIGN.CENTER)

    data_table(slide, [
        ["Metric", "Train", "Test"],
        ["Accuracy", "0.84", "0.81"],
        ["ROC AUC", "0.91", "0.88"],
        ["Sensitivity", "0.87", "0.83"],
        ["Specificity", "0.82", "0.80"],
        ["F1 score", "0.85", "0.82"],
    ], Inches(4.5), Inches(1.25), Inches(3.6), Inches(2.5), font_size=12)

    add_rounded(slide, Inches(8.5), Inches(1.25), Inches(4.2), Inches(2.5), RGBColor(0xFB, 0xF4, 0xDF), line_color=GOLD)
    kd = add_textbox(slide, Inches(8.75), Inches(1.45), Inches(3.75), Inches(2.1))
    fill_box(kd, [
        ("KEY FEATURE CONFIRMED", 12, True, GOLD),
        ("Phenolic OH on an aromatic ring is strongly associated with anti Salmonella activity.", 15, True, NAVY),
        ("The model reproduces this well known essential oil chemistry, confirming it predicts correctly.", 11, False, LIGHT),
    ])

    plants = [
        ("plants/botanical-oregano-cut.png", "Oregano", "carvacrol", "bits/bit_1607.png", True),
        ("plants/botanical-thyme-cut.png", "Wild thyme", "thymol", "bits/bit_1607.png", True),
        ("plants/botanical-clove-cut.png", "Clove", "eugenol", "bits/bit_1607.png", True),
        ("plants/botanical-blackpepper-cut.png", "Black pepper", "β-pinene", "bits/bit_549.png", False),
    ]
    for i, (pimg, name, mol, bit, ok) in enumerate(plants):
        x = 0.7 + i * 3.15
        add_picture(slide, pimg, Inches(x + 0.45), Inches(4.55), height=Inches(1.25))
        add_picture(slide, bit, Inches(x + 0.7), Inches(5.82), height=Inches(0.62))
        nb = add_textbox(slide, Inches(x), Inches(6.48), Inches(3.0), Inches(0.5))
        label = "Bit 1607 positive" if ok else "Bit 549 negative"
        fill_box(nb, [(name, 13, True, NAVY), (f"{mol} · {label}", 11, False, GREEN if ok else RUST)], align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_notes(slide, "We first encoded the molecules using 2048-bit Morgan fingerprints, which is a standard setting. After removing the fragments that were never present in our dataset, we were left with 682 descriptors.\n\nThen we removed one member of each highly correlated pair (|r| > 0.95), reducing the number of descriptors to 337. After that, we ranked the remaining descriptors by permutation feature importance, keeping only those whose shuffling lowered ROC AUC by at least 0.01. Exactly ten descriptors passed this threshold.\n\nTo evaluate the model, we used repeated stratified five fold cross validation with 100 repeats. The area under the curve was 0.88, meaning the model ranked active oils above inactive ones very well. About 88% of the time, a randomly chosen active oil received a higher predicted probability than a randomly chosen inactive oil. Under Hosmer's scale, this is excellent discrimination, although it remains internal validation on our dataset, not an external holdout set. We also used 1000 bootstrap samples to estimate uncertainty around the regression coefficients, not to validate predictive performance.\n\nThe training accuracy was 0.84 and the test accuracy was 0.81, so the gap was small. That suggests the model generalised well and was not simply memorising the training data. Test sensitivity was 0.83 and specificity 0.80.\n\nBased on this, we could interpret the selected features more confidently. Phenolic structures, especially those encoded by Bit_1607, were strongly associated with activity against Salmonella, while branched bicyclic hydrocarbons, such as Bit_549, were associated with inactivity. This is an associative model finding, not a claim of molecular causation. Partial dependence analysis showed that when strong positive bits were absent the predicted probability was about 0.45, and when they were present it rose to as high as about 0.95.\n\nInterestingly, none of the four essential oils we tested experimentally was rich in phenolic compounds. Instead, they were dominated by oxygenated monoterpenes. This helps explain why their activity was modest compared with the phenolic rich oils reported by Soković and colleagues, who documented that carvacrol and thymol containing oils, which differ only in hydroxyl orientation on the benzene ring, showed MICs as low as 0.5 µg/mL against S. Typhimurium.\n\nCarvacrol and thymol are signature compounds of plants such as wild thyme and oregano, which are common in Serbia and traditionally used in foods and medicines. Similar antibacterial effects have also been reported for eugenol, another phenolic compound found predominantly in clove oil, which acts on the bacterial membrane and can induce oxidative stress through reactive oxygen species.\n\nThis supports our logistic regression analysis, in which bits encoding an aromatic hydroxyl group receive the strongest positive weights. It is also a reminder of why protecting and studying our local flora is important.")


def slide_limitations(prs):
    slide = blank_slide(prs)
    add_title(slide, "Limitations")

    left = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.7), Inches(5.0))
    fill_box(left, [
        ("Experimental", 20, True, RUST),
        ("1. Only four commercial oils", 15, True, NAVY),
        ("Limited chemotypic coverage; one production batch each.", 13, False, GRAY),
        ("2. Encapsulation was morphological", 15, True, NAVY),
        ("No release kinetics, retained bioactivity, or food matrix tests.", 13, False, GRAY),
        ("3. Single pathogen strain", 15, True, NAVY),
        ("Only S. Typhimurium ATCC 14028; no clinical isolates or biofilms.", 13, False, GRAY),
    ])
    right = add_textbox(slide, Inches(6.9), Inches(1.5), Inches(5.7), Inches(5.0))
    fill_box(right, [
        ("Computational", 20, True, NAVY),
        ("1. Heterogeneous literature data", 15, True, NAVY),
        ("MIC protocols, media, and strains vary across sources.", 13, False, GRAY),
        ("2. Internal validation only", 15, True, NAVY),
        ("Repeated stratified CV; no external holdout set.", 13, False, GRAY),
        ("3. Candidates, not confirmed mechanisms", 15, True, NAVY),
        ("Highlighted substructures still need wet lab confirmation.", 13, False, GRAY),
    ])
    add_footer(slide)
    add_notes(slide, "Like every study, this work has limitations that should be considered. In the experimental phase, we tested only four commercial essential oils and one Salmonella Typhimurium strain. We focused mainly on bead morphology and did not evaluate release kinetics, biological activity after encapsulation, or performance in food systems.\n\nIn the computational phase, the dataset came from published studies that used different experimental protocols, and only major compounds present above 10% were included. Although the model performed well, it was evaluated only by internal repeated stratified cross validation, with bootstrapping used to estimate coefficient uncertainty rather than predictive performance. The identified molecular features should therefore be viewed as promising candidates that require confirmation in future laboratory experiments.")


def slide_future(prs):
    slide = blank_slide(prs)
    add_title(slide, "Future Directions")
    sub = add_textbox(slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.35))
    write_text(sub.text_frame, "From broader wet lab work to a closed AI to lab loop", size=14, color=LIGHT)

    left = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.7), Inches(3.6))
    fill_box(left, [
        ("Experimental", 18, True, GREEN),
        ("1. More oils & microorganisms", 14, True, NAVY),
        ("Analyse a larger set of essential oils and microbes", 12, False, GRAY),
        ("2. Different encapsulation systems", 14, True, NAVY),
        ("Explore alternative carriers beyond alginate beads", 12, False, GRAY),
        ("3. Real food applications", 14, True, NAVY),
        ("Test formulations in actual food systems", 12, False, GRAY),
    ])
    right = add_textbox(slide, Inches(6.9), Inches(1.5), Inches(5.7), Inches(3.6))
    fill_box(right, [
        ("Computational", 18, True, NAVY),
        ("1. Extend the pipeline", 14, True, NAVY),
        ("Apply the same approach with modern AI methods", 12, False, GRAY),
        ("2. Graph nets & molecular LMs", 14, True, NAVY),
        ("Richer models of chemical structure and activity", 12, False, GRAY),
        ("3. Active learning", 14, True, NAVY),
        ("Predictions guide the next laboratory experiments", 12, False, GRAY),
    ])
    add_rounded(slide, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.35), GREEN_TINT, line_color=GOLD)
    loop = add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.1))
    fill_box(loop, [
        ("Predict  →  Test  →  Retrain  →  repeat", 16, True, GOLD),
        ("Goal: continuous AI and lab feedback for faster discovery", 14, True, NAVY),
    ])
    add_footer(slide)
    add_notes(slide, "This work opens several directions for future research, both experimentally and computationally.\n\nOn the experimental side, we would like to analyse a larger number of essential oils and microorganisms, explore different encapsulation systems, and test our formulations in real food applications.\n\nOn the computational side, the same pipeline can be extended using modern AI methods such as graph neural networks, molecular language models, and active learning, where predictions guide the next laboratory experiments.\n\nUltimately, the goal is to create a continuous feedback loop between AI and laboratory validation, making the discovery of new plant derived antimicrobials faster and more efficient.")


def slide_thanks(prs):
    slide = blank_slide(prs)
    box = add_textbox(slide, Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.2))
    write_text(box.text_frame, "Thank you", size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    sub = add_textbox(slide, Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.6))
    write_text(sub.text_frame, "I welcome your questions", size=20, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "That concludes my presentation. I hope I have shown how combining experimental work with machine learning can help us better understand what drives the antimicrobial potential of plant essential oils and guide future discoveries of active molecular sites.\n\nThank you very much for your attention, and I will be happy to answer any questions.")


def build():
    prs = Presentation()
    set_slide_size(prs)

    slide_title(prs)
    slide_acknowledgments(prs)
    slide_roadmap(prs)
    slide_why_hub(prs)
    slide_lab_hub(prs)
    slide_gcms(prs)
    slide_bioactivity(prs)
    slide_encap(prs)
    slide_ml_pipeline(prs)
    slide_findings(prs)
    slide_limitations(prs)
    slide_future(prs)
    slide_thanks(prs)

    for slide in prs.slides:
        set_fade_transition(slide)

    out = ROOT / "PhD_Defence_Presentation_BACKUP.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"Slides (incl. hidden zoom): {len(prs.slides)}")
    return out


if __name__ == "__main__":
    build()
