import re
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn


ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "img"

NAVY = RGBColor(0x1A, 0x3A, 0x5C)
GREEN = RGBColor(0x4A, 0x67, 0x41)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
RUST = RGBColor(0xB1, 0x35, 0x07)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, size=18, bold=False, color=NAVY, italic=False, font_name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


ITALIC_RE = re.compile(r"(Salmonella|S\.(?= Typhimurium))")


def emit_runs(paragraph, text, size, bold, color, base_italic=False):
    parts = ITALIC_RE.split(text)
    for i, part in enumerate(parts):
        if part == "":
            continue
        run = paragraph.add_run()
        run.text = part
        set_run(run, size=size, bold=bold, color=color, italic=base_italic or i % 2 == 1)


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def write_text(tf, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, italic=False, clear=True):
    if clear:
        tf.clear()
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    emit_runs(p, text, size, bold, color, base_italic=italic)
    return p


def add_para(tf, text, size=16, bold=False, color=GRAY, align=PP_ALIGN.LEFT, italic=False, space_before=4, space_after=2):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    emit_runs(p, text, size, bold, color, base_italic=italic)
    return p


def add_title(slide, text, top=0.35, size=32, color=NAVY):
    box = add_textbox(slide, Inches(0.6), Inches(top), Inches(12.1), Inches(0.7))
    write_text(box.text_frame, text, size=size, bold=True, color=color)
    return box


def add_kicker(slide, text, color=GREEN, top=0.18):
    box = add_textbox(slide, Inches(0.6), Inches(top), Inches(12.1), Inches(0.35))
    write_text(box.text_frame, text.upper(), size=14, bold=True, color=color)
    return box


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text.replace("<br>", "\n").replace("<br/>", "\n").replace("<i>", "").replace("</i>", "").replace("<strong>", "").replace("</strong>", "")


def prepare_image(path):
    path = Path(path)
    if not path.exists():
        return None
    if path.suffix.lower() in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"}:
        return path
    from PIL import Image
    cache = ROOT / "scripts" / "_pptx_img_cache"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{path.stem}.png"
    if not out.exists() or out.stat().st_mtime < path.stat().st_mtime:
        Image.open(path).convert("RGBA").save(out, "PNG")
    return out


def add_picture(slide, rel_path, left, top, width=None, height=None):
    path = prepare_image(IMG / rel_path)
    if path is None:
        return None
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def crop_square(path):
    from PIL import Image
    path = prepare_image(path)
    if path is None:
        return None
    im = Image.open(path).convert("RGB")
    w, h = im.size
    s = min(w, h)
    left = (w - s) // 2
    top = (h - s) // 2
    im = im.crop((left, top, left + s, top + s))
    max_px = 500
    if s > max_px:
        im = im.resize((max_px, max_px), Image.LANCZOS)
    cache = ROOT / "scripts" / "_pptx_img_cache"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{Path(path).stem}_sq.jpg"
    im.save(out, "JPEG", quality=85)
    return out


def add_circular_picture(slide, rel_path, left, top, diameter, ring=None, ring_pt=1.5):
    square = crop_square(IMG / rel_path)
    if square is None:
        return None
    pic = slide.shapes.add_picture(str(square), left, top, width=diameter, height=diameter)
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    geom = parse_xml(
        '<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'prst="ellipse"><a:avLst/></a:prstGeom>'
    )
    spPr.append(geom)
    if ring is not None:
        pic.line.color.rgb = ring
        pic.line.width = Pt(ring_pt)
    return pic


def set_fade_transition(slide, speed="slow"):
    sld = slide._element
    old = sld.find(qn("p:transition"))
    if old is not None:
        sld.remove(old)
    trans = parse_xml(
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'spd="%s"><p:fade/></p:transition>' % speed
    )
    clr = sld.find(qn("p:clrMapOvr"))
    if clr is not None:
        clr.addnext(trans)
    else:
        sld.find(qn("p:cSld")).addnext(trans)


def add_footer(slide, tagline="Bridging phytochemistry & AI"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(9.0), Inches(0.33))
    write_text(box.text_frame, tagline, size=10, color=GOLD)
    return box


def add_zoomable_picture(prs, slide, rel_path, left, top, width=None, height=None):
    from PIL import Image
    thumb = add_picture(slide, rel_path, left, top, width=width, height=height)
    if thumb is None:
        return None
    path = prepare_image(IMG / rel_path)
    iw, ih = Image.open(path).size

    zoom = blank_slide(prs)
    zoom._element.set("show", "0")
    backdrop = zoom.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    backdrop.fill.solid()
    backdrop.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    backdrop.line.fill.background()

    margin = Inches(0.3)
    avail_w = SLIDE_W - 2 * margin
    avail_h = SLIDE_H - 2 * margin
    scale = min(avail_w / iw, avail_h / ih)
    bw = int(iw * scale)
    bh = int(ih * scale)
    big = zoom.shapes.add_picture(str(path), int((SLIDE_W - bw) / 2), int((SLIDE_H - bh) / 2),
                                  width=bw, height=bh)
    hint = zoom.shapes.add_textbox(Inches(0.4), Inches(7.08), Inches(6.0), Inches(0.32))
    write_text(hint.text_frame, "Click image to go back", size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

    thumb.click_action.target_slide = zoom
    backdrop.click_action.target_slide = slide
    big.click_action.target_slide = slide
    return thumb


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_size(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
